"""Tests for app.utils.pptx_converter.

- extract_slide_notes: tested with an in-memory generated PPTX (no LibreOffice needed).
- decide_narration_source: pure logic, always runs.
- pptx_to_images: smoke test; skipped when soffice/pdftoppm are not on PATH.
"""

import os
import shutil
import tempfile
import unittest
from pathlib import Path

from app.utils import pptx_converter


def _make_test_pptx(path: Path, slides_with_notes: list[str | None]) -> None:
    """Build a minimal .pptx in-memory with the given notes per slide.

    Uses python-pptx; no system tools required.
    """
    from pptx import Presentation

    prs = Presentation()
    for note in slides_with_notes:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        if note is not None:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = note
    prs.save(str(path))


class TestDecideNarrationSource(unittest.TestCase):
    def test_all_with_notes(self):
        self.assertEqual(
            pptx_converter.decide_narration_source(["a", "b", "c"]),
            [True, True, True],
        )

    def test_all_without_notes(self):
        self.assertEqual(
            pptx_converter.decide_narration_source([None, None]),
            [False, False],
        )

    def test_mixed(self):
        self.assertEqual(
            pptx_converter.decide_narration_source(["x", None, "y", None]),
            [True, False, True, False],
        )

    def test_empty_input(self):
        self.assertEqual(pptx_converter.decide_narration_source([]), [])


class TestExtractSlideNotes(unittest.TestCase):
    def setUp(self) -> None:
        self.tmp = tempfile.mkdtemp(prefix="mpt_pptx_test_")

    def tearDown(self) -> None:
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_returns_none_for_slides_without_notes(self):
        pptx_path = Path(self.tmp) / "no_notes.pptx"
        _make_test_pptx(pptx_path, [None, None, None])
        result = pptx_converter.extract_slide_notes(pptx_path)
        self.assertEqual(result, [None, None, None])

    def test_returns_text_for_slides_with_notes(self):
        pptx_path = Path(self.tmp) / "with_notes.pptx"
        _make_test_pptx(pptx_path, ["第一页备注", "第二页备注"])
        result = pptx_converter.extract_slide_notes(pptx_path)
        self.assertEqual(result, ["第一页备注", "第二页备注"])

    def test_mixed_notes_returns_none_for_empty(self):
        pptx_path = Path(self.tmp) / "mixed.pptx"
        _make_test_pptx(pptx_path, ["有备注", None, "  ", "又有备注"])
        # "  "（仅空白）应视为 None；非空文本保留 strip 结果
        result = pptx_converter.extract_slide_notes(pptx_path)
        self.assertEqual(result, ["有备注", None, None, "又有备注"])

    def test_length_matches_slide_count(self):
        pptx_path = Path(self.tmp) / "five.pptx"
        _make_test_pptx(pptx_path, [None] * 5)
        result = pptx_converter.extract_slide_notes(pptx_path)
        self.assertEqual(len(result), 5)

    def test_missing_file_raises(self):
        with self.assertRaises(pptx_converter.PPTXConversionError):
            pptx_converter.extract_slide_notes(Path(self.tmp) / "nope.pptx")


class TestPptxToImages(unittest.TestCase):
    """Smoke test. Requires soffice + pdftoppm on PATH (e.g. Docker image)."""

    def setUp(self) -> None:
        self.tmp = tempfile.mkdtemp(prefix="mpt_pptx_render_")

    def tearDown(self) -> None:
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_smoke_renders_each_slide(self):
        if not (shutil.which("soffice") or shutil.which("libreoffice")):
            self.skipTest("soffice not on PATH")
        if not shutil.which("pdftoppm"):
            self.skipTest("pdftoppm not on PATH")

        pptx_path = Path(self.tmp) / "deck.pptx"
        _make_test_pptx(pptx_path, ["n1", None, "n3"])
        out_dir = Path(self.tmp) / "out"
        images = pptx_converter.pptx_to_images(pptx_path, out_dir)
        self.assertEqual(len(images), 3)
        for p in images:
            self.assertTrue(Path(p).is_file())
            self.assertTrue(p.lower().endswith(".png"))


if __name__ == "__main__":
    unittest.main()
