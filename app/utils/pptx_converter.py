"""
PPTX → 每页图片 + 备注提取（webui「PPT 转视频」页用）。

渲染：LibreOffice headless（`soffice`）把 .pptx 转 PDF，再用 `pdftoppm`（poppler-utils）把
每页切成 PNG，输出文件名 `slide-1.png ... slide-N.png`（pdftoppm 默认零填充）。
备注：python-pptx 读取每页的 notes_slide 文本，纯 Python、无系统依赖。

依赖（系统）：libreoffice-impress、poppler-utils、fonts-noto-cjk（Dockerfile 已配）。
依赖（Python）：python-pptx。
"""

from __future__ import annotations

import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

# pdftoppm 输出的文件名前缀与零填充宽度
_PAGE_PREFIX = "slide"
_PAGE_PADDING = 3  # slide-001.png

# 渲染 / 转换超时（秒）。20 页大 PPT 通常 < 30s，留 180s 余量。
_TIMEOUT_SEC = 180

# pdftoppm 分辨率（DPI）。150 是经验值，字够清、文件不大。
_DPI = 150

# LibreOffice 用户配置文件子目录（避免容器内多用户冲突）
_LO_PROFILE_DIRNAME = ".lo_profile"


class PPTXConversionError(RuntimeError):
    """PPTX 转换失败时抛出。"""


def _find_soffice() -> str:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise PPTXConversionError(
            "未找到 LibreOffice（soffice）。请在系统中安装 libreoffice-impress，"
            "或使用 Docker 镜像。Windows 本地可改用 PowerPoint COM 方案。"
        )
    return soffice


def pptx_to_images(pptx_path: str | os.PathLike, out_dir: str | os.PathLike) -> list[str]:
    """把 .pptx 转成每页一张 PNG，按页码排序返回绝对路径。

    Args:
        pptx_path: 输入 .pptx 绝对路径。
        out_dir: 输出目录（不存在会自动创建）。

    Returns:
        排序后的 PNG 绝对路径列表，长度等于幻灯片页数。

    Raises:
        PPTXConversionError: soffice 缺失、转换失败或 PDF→PNG 失败。
    """
    pptx_path = Path(pptx_path).resolve()
    out_dir = Path(out_dir).resolve()
    if not pptx_path.is_file():
        raise PPTXConversionError(f"PPTX 文件不存在: {pptx_path}")
    out_dir.mkdir(parents=True, exist_ok=True)

    soffice = _find_soffice()

    # 1) soffice → PDF。配独立 user profile（避免容器内 root profile 锁问题），把 HOME 指向临时
    # 目录以防 LibreOffice 试图写 ~/.config。
    with tempfile.TemporaryDirectory(prefix="mpt_lo_") as tmp_home:
        env = os.environ.copy()
        env["HOME"] = tmp_home
        profile_dir = out_dir / _LO_PROFILE_DIRNAME
        profile_dir.mkdir(exist_ok=True)
        cmd = [
            soffice,
            "--headless",
            "--norestore",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            f"-env:UserInstallation=file://{profile_dir}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(pptx_path),
        ]
        proc = subprocess.run(
            cmd,
            env=env,
            capture_output=True,
            text=True,
            timeout=_TIMEOUT_SEC,
        )
        if proc.returncode != 0:
            raise PPTXConversionError(
                f"LibreOffice 转换失败 (rc={proc.returncode}): "
                f"{proc.stderr.strip() or proc.stdout.strip()}"
            )

    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.is_file():
        raise PPTXConversionError(f"LibreOffice 未生成 PDF: {pdf_path}")

    # 2) pdftoppm → 每页 PNG
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        raise PPTXConversionError("未找到 pdftoppm（poppler-utils）。")
    proc = subprocess.run(
        [
            pdftoppm,
            "-png",
            "-r",
            str(_DPI),
            str(pdf_path),
            str(out_dir / _PAGE_PREFIX),
        ],
        capture_output=True,
        text=True,
        timeout=_TIMEOUT_SEC,
    )
    if proc.returncode != 0:
        raise PPTXConversionError(
            f"pdftoppm 切图失败 (rc={proc.returncode}): "
            f"{proc.stderr.strip() or proc.stdout.strip()}"
        )

    # 清理中间产物
    try:
        pdf_path.unlink()
    except OSError:
        pass

    images = sorted(
        (out_dir / _PAGE_PREFIX).parent.glob(f"{_PAGE_PREFIX}-*.png"),
        key=_natural_page_key,
    )
    if not images:
        raise PPTXConversionError("pdftoppm 未产出任何 PNG")
    return [str(p.resolve()) for p in images]


def _natural_page_key(path: Path) -> int:
    """从 slide-001.png / slide-1.png 里抽页码，整数排序。"""
    m = re.search(r"-(\d+)\.png$", path.name)
    return int(m.group(1)) if m else 0


def extract_slide_notes(pptx_path: str | os.PathLike) -> list[str | None]:
    """用 python-pptx 读每页备注。空备注返回 None，非空返回 strip 后的文本。

    返回列表长度 == 幻灯片页数，与 pptx_to_images 的输出一一对应。
    """
    from pptx import Presentation  # 局部 import，避免在系统无 python-pptx 时影响其他模块

    pptx_path = Path(pptx_path).resolve()
    if not pptx_path.is_file():
        raise PPTXConversionError(f"PPTX 文件不存在: {pptx_path}")

    prs = Presentation(str(pptx_path))
    notes: list[str | None] = []
    for slide in prs.slides:
        text: str | None = None
        if slide.has_notes_slide:
            raw = slide.notes_slide.notes_text_frame.text or ""
            stripped = raw.strip()
            if stripped:
                text = stripped
        notes.append(text)
    return notes


def decide_narration_source(notes: list[str | None]) -> list[bool]:
    """按页决定文案来源：True = 用备注，False = 走视觉识别。"""
    return [n is not None for n in notes]
